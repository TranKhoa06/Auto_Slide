import os
import json
import uuid
import asyncio
import copy
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from dotenv import load_dotenv
from google import genai
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches
import fitz
import requests
import urllib.parse

load_dotenv()

app = FastAPI()

# Cấu hình CORS để Frontend gọi được API
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Phục vụ các file tĩnh (Giao diện web)
app.mount("/static", StaticFiles(directory="static"), name="static")

# Thư mục tạm để lưu file
TEMP_DIR = "temp"
os.makedirs(TEMP_DIR, exist_ok=True)

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return "\n".join(text)

async def generate_script_content(file_path: str, slide_count: int, api_key: str):
    """BƯỚC 1: Gọi AI đọc PDF/PPTX và lập Kịch bản JSON"""
    if not api_key:
        api_key = os.environ.get("GEMINI_API_KEY", "")
    client = genai.Client(api_key=api_key)
    print(f"Uploading/Processing {file_path} to Gemini for script generation...")
    
    script_prompt = f"""Bạn là chuyên gia thiết kế nội dung bài giảng. Dựa vào tài liệu học thuật đính kèm, hãy tạo Kịch bản Thuyết trình gồm đúng {slide_count} Slide.
YÊU CẦU NGHIÊM NGẶT:
- GIỮ NGUYÊN NGÔN NGỮ CỦA TÀI LIỆU GỐC (Ví dụ: Nếu tài liệu tiếng Anh, slide phải là tiếng Anh. Tuyệt đối không tự động dịch).
- Dữ liệu trả về PHẢI là định dạng JSON mảng hợp lệ. KHÔNG bọc bằng markdown json.
- Mỗi slide chứa 3 trường: "title" (tiêu đề ngắn gọn < 30 ký tự), "bullets" (mảng gồm 3-5 gạch đầu dòng, mỗi gạch đầu dòng KHÔNG QUÁ 100 ký tự), "speaker_notes" (kịch bản chi tiết để người thuyết trình đọc).
- Cố gắng giữ lại các định nghĩa, công thức, số liệu quan trọng nhất. Đưa các diễn giải dài dòng vào speaker_notes.

MẪU JSON TRẢ VỀ (chỉ trả về đoạn JSON này, không giải thích thêm):
[
  {{
    "title": "Tổng quan AI",
    "bullets": [
      "AI là trí tuệ nhân tạo.",
      "Giúp tự động hóa công việc."
    ],
    "speaker_notes": "Xin chào mọi người..."
  }}
]"""
    print("Generating script...")
    
    if file_path.lower().endswith('.pptx'):
        extracted_text = extract_text_from_pptx(file_path)
        contents = [f"Đây là nội dung trích xuất từ Slide cũ:\n\n{extracted_text}", script_prompt]
        response = client.models.generate_content(
            model='gemini-flash-latest',
            contents=contents
        )
    else:
        sample_file = client.files.upload(file=file_path, config={'display_name': "Lecture Document"})
        response = client.models.generate_content(
            model='gemini-flash-latest',
            contents=[sample_file, script_prompt]
        )
        client.files.delete(name=sample_file.name)
    
    text = response.text.strip()
    if text.startswith("```json"): text = text[7:-3].strip()
    elif text.startswith("```"): text = text[3:-3].strip()
    return text

def apply_dark_mode_theme(slide):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    
    background = slide.background
    fill = background.fill
    fill.solid()
    # Dark Navy / Charcoal background
    fill.fore_color.rgb = RGBColor(20, 25, 35)
    
    # Top neon cyan thin accent line
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0, 255, 255) # Cyan
    shape1.line.fill.background()
    
    # Bottom neon orange thin accent line
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.42), Inches(10), Inches(0.08))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(255, 100, 0) # Orange
    shape2.line.fill.background()
    
    # A small tech-style geometric block in the top left
    shape3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0), Inches(1.5), Inches(0.12))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(255, 100, 0) # Orange
    shape3.line.fill.background()

def set_font(run, size_pt, bold=False, rgb=(50, 50, 50)):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    run.font.name = "Segoe UI"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*rgb)

def generate_native_pptx(json_script, output_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    try:
        script_data = json.loads(json_script)
    except Exception as e:
        raise ValueError(f"JSON kịch bản không hợp lệ: {e}")

    for slide_data in script_data:
        slide = prs.slides.add_slide(blank_layout)
        apply_dark_mode_theme(slide)
        
        # Add Title (Width reduced to 7.5 to avoid FPT logo on the right)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(7.5), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        # White bold title
        for run in p.runs:
            set_font(run, 40, bold=True, rgb=(255, 255, 255))
        
        # Add FPT Logo
        logo_path = "logo_fpt.png"
        import os
        if os.path.exists(logo_path):
            try:
                slide.shapes.add_picture(logo_path, Inches(8.2), Inches(0.3), width=Inches(1.5))
            except Exception:
                pass
        
        # Add Body (Bullets)
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.5) if slide_data.get("image_url") else Inches(9), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        bullets = slide_data.get("bullets", [])
        if isinstance(bullets, list):
            bullets = "\n".join(bullets)
        
        paragraphs = bullets.split('\n')
        for i, para_text in enumerate(paragraphs):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = para_text
            for run in p.runs:
                # Light grey body text for contrast
                set_font(run, 22, rgb=(220, 225, 230))
        
        # Add Speaker Notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes_slide.notes_text_frame.text = slide_data.get("speaker_notes", "")

        # Chèn ảnh vào slide nếu có
        img_url = slide_data.get("image_url")
        if img_url:
            try:
                import urllib.request
                img_path = os.path.join(os.path.dirname(output_path), f"temp_img_{uuid.uuid4().hex[:6]}.jpg")
                if img_url.startswith("http"):
                    headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
                    response = requests.get(img_url, headers=headers, timeout=10)
                    response.raise_for_status()
                    with open(img_path, 'wb') as f:
                        f.write(response.content)
                else:
                    img_path = img_url
                
                slide.shapes.add_picture(img_path, Inches(6.0), Inches(1.8), width=Inches(3.5))
                
            except Exception as e:
                print(f"Error inserting image: {repr(e)}")

    prs.save(output_path)


@app.post("/generate_script")
async def api_generate_script(
    source_file: UploadFile = File(...),
    slide_count: int = Form(...),
    api_key: str = Form("")
):
    try:
        session_id = str(uuid.uuid4())
        ext = ".pdf"
        if source_file.filename and source_file.filename.lower().endswith(".pptx"):
            ext = ".pptx"
            
        file_path = os.path.join(TEMP_DIR, f"{session_id}{ext}")
        with open(file_path, "wb") as f:
            f.write(await source_file.read())
            
        script_text = await generate_script_content(file_path, slide_count, api_key)
        return {"session_id": session_id, "script": script_text, "pdf_path": file_path}
    except Exception as e:
        print("Script Error:", str(e))
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate_images")
async def api_generate_images(
    session_id: str = Form(...),
    script: str = Form(...),
    pdf_path: str = Form(...)
):
    """BƯỚC 1.5: Trích xuất ảnh hoặc sinh ảnh AI cho các slide"""
    try:
        script_data = json.loads(script)
        
        # Thử trích xuất ảnh từ PDF (bỏ qua nếu là PPTX)
        extracted_images = []
        if pdf_path.lower().endswith(".pdf"):
            try:
                doc = fitz.open(pdf_path)
                for page_num in range(min(len(doc), 10)): # Quét 10 trang đầu
                    for img in doc.get_page_images(page_num):
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        img_path = os.path.join(TEMP_DIR, f"{session_id}_img_{len(extracted_images)}.png")
                        if pix.n - pix.alpha < 4:       # this is GRAY or RGB
                            pix.save(img_path)
                        else:               # CMYK: convert to RGB first
                            pix1 = fitz.Pixmap(fitz.csRGB, pix)
                            pix1.save(img_path)
                            pix1 = None
                        pix = None
                        extracted_images.append(img_path)
            except Exception as e:
                print("Lỗi trích xuất ảnh:", e)

        # Gán ảnh cho script
        for i, slide in enumerate(script_data):
            if i < len(extracted_images):
                slide["image_url"] = extracted_images[i]
                slide["image_source"] = "extracted"
            else:
                # Sinh ảnh bằng AI (Pollinations.ai)
                prompt = urllib.parse.quote(f"Minimalist illustration for presentation slide about {slide.get('title')}. Modern corporate style, clean background, highly detailed")
                img_url = f"https://image.pollinations.ai/prompt/{prompt}?width=800&height=600&nologo=true"
                slide["image_url"] = img_url
                slide["image_source"] = "ai_generated"

        return {"script": json.dumps(script_data, ensure_ascii=False, indent=2)}
    except Exception as e:
        print("Image Gen Error:", str(e))
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate_pptx")
async def api_generate_pptx(
    session_id: str = Form(...),
    script: str = Form(...)
):
    try:
        output_path = os.path.join(TEMP_DIR, f"{session_id}_output.pptx")
        generate_native_pptx(script, output_path)
        
        return FileResponse(
            path=output_path, 
            filename="AutoSlide_Generated.pptx", 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        print("PPTX Error:", str(e))
        raise HTTPException(status_code=500, detail=str(e))

# Phục vụ Frontend HTML
app.mount("/", StaticFiles(directory="static", html=True), name="static")
