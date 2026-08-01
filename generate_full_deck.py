import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

def extract_image_from_pdf(pdf_path, page_num, output_img):
    doc = fitz.open(pdf_path)
    page = doc[page_num]
    image_list = page.get_images(full=True)
    if image_list:
        largest_xref = None
        max_size = 0
        for img in image_list:
            xref = img[0]
            base_image = doc.extract_image(xref)
            size = len(base_image["image"])
            if size > max_size:
                max_size = size
                largest_xref = xref
                
        if largest_xref:
            base_image = doc.extract_image(largest_xref)
            image_bytes = base_image["image"]
            with open(output_img, "wb") as f:
                f.write(image_bytes)
            return True
    return False

def apply_dark_mode_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 25, 35) # Dark Navy/Charcoal
    
    # Top cyan line
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0, 255, 255)
    shape1.line.fill.background()
    
    # Bottom orange line
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.42), Inches(10), Inches(0.08))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(255, 100, 0)
    shape2.line.fill.background()
    
    # Top left small accent block
    shape3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0), Inches(1.5), Inches(0.12))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(255, 100, 0)
    shape3.line.fill.background()
    
    # Add FPT Logo
    logo_path = "logo_fpt.png"
    if os.path.exists(logo_path):
        try:
            # Góc trên bên phải, xích qua 8.2 để không đè title
            slide.shapes.add_picture(logo_path, Inches(8.2), Inches(0.3), width=Inches(1.5))
        except Exception:
            pass

def set_font(run, size_pt, bold=False, rgb=(50, 50, 50)):
    run.font.name = "Segoe UI"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*rgb)

def add_title(slide, text):
    # Giới hạn width=7.5 để không đè vào Logo FPT
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(7.5), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        set_font(run, 40, bold=True, rgb=(255, 255, 255))

def add_body(slide, text, top_inch=1.5, width_inch=9):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top_inch), Inches(width_inch), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    paragraphs = text.split('\n')
    for i, para_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para_text
        for run in p.runs:
            set_font(run, 22, rgb=(220, 225, 230))

def main():
    pdf_path = "Chapter 11-Simple Linear Regression and Correlation.pdf"
    
    # Define images to extract
    image_map = {
        "table_11_1": {"page": 3, "file": "img_table.png"},
        "scatter_1": {"page": 4, "file": "img_scatter1.png"},
        "deviations": {"page": 8, "file": "img_deviations.png"},
        "fitted_scatter": {"page": 14, "file": "img_fitted.png"},
        "conf_limits_scatter": {"page": 29, "file": "img_conf.png"},
        "corr_examples": {"page": 35, "file": "img_corr.png"}
    }
    
    print("Extracting images from PDF...")
    for key, data in image_map.items():
        extract_image_from_pdf(pdf_path, data["page"], data["file"])

    # Full content mapping
    ai_json = [
        {
            "layout": "title_slide",
            "title": "Simple Linear Regression\nand Correlation",
            "subtitle": "Chapter 11 - Probability & Statistics"
        },
        {
            "layout": "one_column",
            "title": "Learning Objectives",
            "content": "1. Empirical Models\n2. Simple Linear Regression\n3. Properties of the Least Squares Estimators\n4. Hypothesis Tests in Simple Linear Regression\n5. Correlation"
        },
        {
            "layout": "one_column",
            "title": "1. Empirical Models",
            "content": "• Many problems in engineering and science involve exploring relationships between two or more variables.\n• Regression analysis is a statistical technique used to build models to predict outcomes.\n• Example: Predicting chemical product yield based on operating temperature."
        },
        {
            "layout": "image_slide",
            "title": "Data Example: Oxygen and Hydrocarbon",
            "content": "Table 11-1 shows the observation numbers, hydrocarbon level x(%), and Purity y(%).",
            "image_file": image_map["table_11_1"]["file"]
        },
        {
            "layout": "image_slide",
            "title": "Scatter Diagram",
            "content": "Scatter Diagram of oxygen purity versus hydrocarbon level (Table 11-1). Visualizing the data is always the first step in regression analysis.",
            "image_file": image_map["scatter_1"]["file"]
        },
        {
            "layout": "one_column",
            "title": "2. Simple Linear Regression Model",
            "content": "Based on the scatter diagram, we assume the mean of random variable Y is related to x by a straight-line relationship:\n\nE(Y|x) = β0 + β1x\n\nThe simple linear regression model with random error (ε) is given by:\nY = β0 + β1x + ε\n\nWhere β0 and β1 are regression coefficients."
        },
        {
            "layout": "image_slide",
            "title": "Deviations of the Data",
            "content": "The actual observed values deviate from the estimated regression line. The difference is the error (ε).",
            "image_file": image_map["deviations"]["file"]
        },
        {
            "layout": "one_column",
            "title": "Method of Least Squares",
            "content": "The method of least squares is used to estimate the parameters β0 and β1 by minimizing the sum of the squares of the vertical deviations.\n\nThe estimated regression line is:\nŷ = β̂0 + β̂1x\n\nWhere β̂1 = Sxy / Sxx\nAnd β̂0 = ȳ - β̂1x̄"
        },
        {
            "layout": "image_slide",
            "title": "Fitted Regression Line Example",
            "content": "The fitted simple linear regression model for the Oxygen Purity data is:\n\nŷ = 74.283 + 14.947x",
            "image_file": image_map["fitted_scatter"]["file"]
        },
        {
            "layout": "one_column",
            "title": "3. Properties of the Least Squares Estimator",
            "content": "• The estimator is unbiased: E(β̂1) = β1\n• The error sum of squares (SSE) helps in estimating the variance σ².\n• An unbiased estimator of σ² is:\n\nσ̂² = SSE / (n - 2)"
        },
        {
            "layout": "one_column",
            "title": "4. Hypothesis Tests in Regression",
            "content": "We often need to test if there is a significant linear relationship between x and Y.\n\nTest on the Slope (β1):\nH0: β1 = 0 (No linear relationship)\nH1: β1 ≠ 0 (Significant linear relationship)\n\nWe use a t-statistic: T0 = β̂1 / se(β̂1) which follows a t-distribution with n-2 degrees of freedom."
        },
        {
            "layout": "image_slide",
            "title": "Confidence Intervals on Mean Response",
            "content": "We can construct a 95% confidence interval about the mean response for the data. The limits curve away from the regression line at the edges.",
            "image_file": image_map["conf_limits_scatter"]["file"]
        },
        {
            "layout": "one_column",
            "title": "5. Correlation",
            "content": "Correlation measures the strength of a linear relationship between two variables.\nThe sample correlation coefficient (R) is computed as:\n\nR = Sxy / sqrt(Sxx * SST)\n\nProperties:\n• -1 ≤ R ≤ 1\n• Value does not change if scales change.\n• Value is not affected by interchanging x and y."
        },
        {
            "layout": "image_slide",
            "title": "Types of Correlation",
            "content": "Different scatter plots showing strong negative, strong positive, weak positive, and nonlinear correlation.",
            "image_file": image_map["corr_examples"]["file"]
        },
        {
            "layout": "title_slide",
            "title": "End of Chapter 11",
            "subtitle": "Generated 100% via Native AI Typography"
        }
    ]
    
    print("Bat dau tao 15 Slide hoan chinh...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    for slide_data in ai_json:
        slide = prs.slides.add_slide(blank_layout)
        apply_dark_mode_theme(slide)
        
        if slide_data["layout"] == "title_slide":
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data["title"]
            for run in p.runs:
                set_font(run, 50, bold=True, rgb=(255, 255, 255))
            
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = slide_data["subtitle"]
            for run in p2.runs:
                set_font(run, 28, rgb=(0, 255, 255))
            
        elif slide_data["layout"] == "one_column":
            add_title(slide, slide_data["title"])
            add_body(slide, slide_data["content"], top_inch=1.8)
            
        elif slide_data["layout"] == "image_slide":
            add_title(slide, slide_data["title"])
            # Đặt Text sang cột trái (width=5.0) để dành không gian cột phải cho ảnh
            add_body(slide, slide_data["content"], top_inch=1.8, width_inch=5.0)
            img_path = slide_data.get("image_file")
            if img_path and os.path.exists(img_path):
                try:
                    from PIL import Image
                    with Image.open(img_path) as img:
                        w, h = img.size
                        aspect = w / h
                    
                    # Giới hạn kích thước ảnh ở nửa phải slide
                    max_w = Inches(4.0)
                    max_h = Inches(4.5)
                    max_aspect = max_w / max_h
                    
                    if aspect > max_aspect:
                        target_w = max_w
                        target_h = max_w / aspect
                    else:
                        target_h = max_h
                        target_w = max_h * aspect
                        
                    left_margin = Inches(5.5) + (Inches(4.0) - target_w) / 2
                    top_margin = Inches(1.8) + (Inches(4.5) - target_h) / 2
                    
                    slide.shapes.add_picture(img_path, left_margin, top_margin, width=target_w, height=target_h)
                except Exception as e:
                    print(f"Error inserting image: {repr(e)}")
                
    output_path = "Full_Deck_Chapter_11_New_Style.pptx"
    prs.save(output_path)
    
    # Cleanup temp images
    for data in image_map.values():
        if os.path.exists(data["file"]):
            os.remove(data["file"])
            
    print(f"Thanh cong! Da tao file hoan chinh {output_path}")

if __name__ == "__main__":
    main()
