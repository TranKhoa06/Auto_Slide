import asyncio
import os
import json
from app import generate_script_content, extract_layouts_from_pptx, map_script_to_layout, inject_content_to_pptx
from pptx import Presentation

async def run_tests():
    api_key = ''
    pdf_path = 'Chapter 11-Simple Linear Regression and Correlation.pdf'
    
    print('=====================================')
    print('1. BẮT ĐẦU TẠO KỊCH BẢN TỪ PDF...')
    print('=====================================')
    script = await generate_script_content(pdf_path, 10, api_key)
    print('\n[KỊCH BẢN ĐÃ TẠO]\n', script[:300], '\n...(Lược bớt)...\n')
    
    templates = [
        'Beige Green Playful Computer Hardware Components Group Project Presentation.pptx',
        'Chapter 8.pptx',
        'Red and Yellow Retro Vintage Circus Trivia Night Game Presentation.pptx',
        'SU26_CSD202_IC2002_Group2_Topic15.pptx',
        'Topic 35.pptx'
    ]
    
    for t in templates:
        print('\n=====================================')
        print(f'2. KIỂM THỬ TEMPLATE: {t}')
        print('=====================================')
        
        try:
            prs = Presentation(t)
            layouts = extract_layouts_from_pptx(prs)
            
            print(f' - Trích xuất thành công {len(layouts)} slide layouts (Zero-Shot Space Mapping).')
            print(' - Gọi AI (Gemini) để gắn kịch bản vào Không gian tọa độ...')
            
            mapping = await map_script_to_layout(script, layouts, api_key)
            print(f' - AI Mapping THÀNH CÔNG. Số slide được gán chữ: {len(mapping)}')
            
            output_path = os.path.join('temp', t.replace('.pptx', '_test_out.pptx'))
            inject_content_to_pptx(t, mapping, output_path)
            
            print(f' - Đã lưu file PPTX đầu ra: {output_path}')
            
            # Đọc lại file đầu ra để kiểm chứng text
            print(' - KẾT QUẢ INJECTION:')
            out_prs = Presentation(output_path)
            for i, slide in enumerate(out_prs.slides):
                if i > 2: break # Chỉ in 3 slide đầu để xem
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t_str = shape.text.strip().replace('\n', ' ')
                        if t_str: texts.append(t_str)
                print(f'   + Slide {i}: {texts[:2]} ...')
        except Exception as e:
            print(f' - LỖI MẠNG HOẶC LOGIC: {e}')
            
    print('\n🎉 HOÀN THÀNH TOÀN BỘ VÒNG LẶP KIỂM THỬ!')

if __name__ == "__main__":
    asyncio.run(run_tests())
