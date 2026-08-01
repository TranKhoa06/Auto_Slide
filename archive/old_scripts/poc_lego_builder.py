import copy
from pptx import Presentation
from pptx.util import Inches, Pt

def clone_shape_on_same_slide(shape, slide):
    """Clones a text shape on the SAME slide to preserve relationships."""
    new_element = copy.deepcopy(shape.element)
    slide.shapes._spTree.append(new_element)
    return slide.shapes[-1]

def set_text(shape, new_text):
    if shape.has_text_frame:
        shape.text = new_text

def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)

def main():
    template_path = "Chapter 8.pptx"
    prs = Presentation(template_path)
    
    # Dữ liệu nội dung mô phỏng từ AI
    ai_json = [
        {
            "target_slide_index": 0,
            "layout": "title_slide",
            "title": "HỒI QUY TUYẾN TÍNH\n(LINEAR REGRESSION)",
            "subtitle": "Phân tích Dữ liệu với AI"
        },
        {
            "target_slide_index": 1,
            "layout": "one_column",
            "title": "KHÁI NIỆM",
            "content": "Hồi quy tuyến tính là một phương pháp thống kê dùng để mô hình hóa mối quan hệ giữa biến phụ thuộc và một hoặc nhiều biến độc lập. Đây là thuật toán cốt lõi trong Machine Learning."
        },
        {
            "target_slide_index": 2,
            "layout": "two_column",
            "title": "SO SÁNH ỨNG DỤNG",
            "col1_title": "Dự báo",
            "col1_text": "Dự đoán xu hướng thị trường chứng khoán, giá nhà đất, doanh thu bán hàng dựa trên dữ liệu lịch sử.",
            "col2_title": "Đánh giá",
            "col2_text": "Xác định mức độ ảnh hưởng của các yếu tố (ví dụ: tiền quảng cáo) lên kết quả (doanh số bán hàng)."
        }
    ]
    
    print("🚀 Bắt đầu sửa đổi trực tiếp trên cấu trúc Canva (Lego In-Place)...")
    
    # Chúng ta sẽ giữ lại 3 slide đầu tiên và xóa các slide sau
    for i in range(len(prs.slides) - 1, 2, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
        
    for slide_data in ai_json:
        idx = slide_data["target_slide_index"]
        slide = prs.slides[idx]
        
        # Tìm các khối Lego trên slide này
        title_lego = None
        body_lego = None
        shapes_to_delete = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.lower()
                # Định vị khối Title
                if "chapter 4" in text or "member" in text or "vision & mission" in text:
                    title_lego = shape
                # Định vị khối Body
                elif "introduction" in text or "trần anh khoa" in text or "lorem ipsum" in text:
                    if body_lego is None:
                        body_lego = shape
                    else:
                        shapes_to_delete.append(shape)
                else:
                    # Rác text (ví dụ: www.reallygreatsite.com) -> Đưa vào danh sách xóa
                    shapes_to_delete.append(shape)
                    
        # Xóa các text rác để slide sạch sẽ (chỉ giữ lại đồ họa Canva)
        for shape in shapes_to_delete:
            remove_shape(shape)
            
        # Bắt đầu lắp ráp lại
        if slide_data["layout"] == "title_slide":
            set_text(title_lego, slide_data["title"])
            title_lego.width = Inches(8)
            
            if body_lego:
                set_text(body_lego, slide_data["subtitle"])
                
        elif slide_data["layout"] == "one_column":
            set_text(title_lego, slide_data["title"])
            if body_lego:
                set_text(body_lego, slide_data["content"])
                body_lego.left = Inches(1)
                body_lego.width = Inches(8)
                body_lego.top = Inches(2)
                
        elif slide_data["layout"] == "two_column":
            set_text(title_lego, slide_data["title"])
            
            if body_lego:
                # Cột 1 Title
                t1 = clone_shape_on_same_slide(title_lego, slide)
                t1.left = Inches(0.5)
                t1.top = Inches(2.0)
                t1.width = Inches(4.5)
                set_text(t1, slide_data["col1_title"])
                for paragraph in t1.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(20)
                        
                # Cột 1 Text
                b1 = clone_shape_on_same_slide(body_lego, slide)
                b1.left = Inches(0.5)
                b1.top = Inches(2.5)
                b1.width = Inches(4.2)
                set_text(b1, slide_data["col1_text"])
                
                # Cột 2 Title
                t2 = clone_shape_on_same_slide(title_lego, slide)
                t2.left = Inches(5.2)
                t2.top = Inches(2.0)
                t2.width = Inches(4.5)
                set_text(t2, slide_data["col2_title"])
                for paragraph in t2.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(20)
                        
                # Cột 2 Text
                b2 = clone_shape_on_same_slide(body_lego, slide)
                b2.left = Inches(5.2)
                b2.top = Inches(2.5)
                b2.width = Inches(4.2)
                set_text(b2, slide_data["col2_text"])
                
                # Xóa body_lego gốc vì đã nhân bản ra 2 cột
                remove_shape(body_lego)
                
    output_path = "Slide_Hoan_Chinh_Lego_Fixed.pptx"
    prs.save(output_path)
    print(f"✅ Thành công! Đã sửa lỗi hình ảnh và tạo file {output_path}")

if __name__ == "__main__":
    main()
