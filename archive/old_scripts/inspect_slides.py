from pptx import Presentation

def inspect_slides(file_path):
    prs = Presentation(file_path)
    print(f"Inspecting '{file_path}'...")
    print(f"Number of slides: {len(prs.slides)}\n")
    
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type
            if shape.has_text_frame:
                text = shape.text.replace("\n", " ")[:50]
                print(f"  Shape {j} (Type: {shape_type}): TEXT = '{text}...'")
            else:
                print(f"  Shape {j} (Type: {shape_type}): (No text)")
        print("")
        if i >= 2: # Just inspect first 3 slides
            break

if __name__ == "__main__":
    inspect_slides("Chapter 8.pptx")
