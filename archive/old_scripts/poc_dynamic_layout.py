import json
from pptx import Presentation

def create_dynamic_presentation():
    # Sử dụng Template mặc định của PowerPoint (chứa sẵn các layout cơ bản)
    prs = Presentation()
    
    # Giả lập 1 kịch bản JSON mà AI Đạo diễn đã tính toán và quyết định Bố cục (Layout)
    # AI tự biết khi nào dùng 1 cột, khi nào dùng 2 cột dựa trên độ dài và tính chất nội dung.
    ai_director_plan = [
        {
            "layout_id": 0, # Title Slide
            "title": "Trí Tuệ Nhân Tạo Trong Giáo Dục",
            "subtitle": "Giải pháp thiết kế Slide động thông minh"
        },
        {
            "layout_id": 1, # Title and Content (1 Cột)
            "title": "1. Khó khăn hiện tại",
            "bullets": [
                "Slide tạo bởi AI thường bị tràn chữ vì chỉ dùng 1 khung tĩnh.",
                "Thiếu tính thẩm mỹ do không biết tự động ngắt trang.",
                "Không tận dụng được khoảng trắng và hệ thống phân cấp thị giác."
            ]
        },
        {
            "layout_id": 3, # Two Content (2 Cột)
            "title": "2. Giải pháp: Dynamic Layout Library",
            "column_1": [
                "Bên trái: Vai trò của AI",
                "Đóng vai trò như một Art Director.",
                "Phân tích độ dài văn bản.",
                "Chỉ định mã Layout phù hợp (1 cột, 2 cột, hoặc quote)."
            ],
            "column_2": [
                "Bên phải: Vai trò của Code",
                "Đóng vai trò người thợ lắp ráp.",
                "Khởi tạo đúng Slide Master theo mã Layout của AI.",
                "Đổ dữ liệu chính xác vào các placeholder tương ứng."
            ]
        }
    ]
    
    print("🚀 Bắt đầu tạo Slide bằng kiến trúc Đạo diễn AI...")
    
    for slide_data in ai_director_plan:
        layout_id = slide_data["layout_id"]
        slide_layout = prs.slide_layouts[layout_id]
        slide = prs.slides.add_slide(slide_layout)
        
        # Bố cục 0: Title Slide
        if layout_id == 0:
            slide.shapes.title.text = slide_data["title"]
            slide.placeholders[1].text = slide_data["subtitle"]
            
        # Bố cục 1: 1 Cột
        elif layout_id == 1:
            slide.shapes.title.text = slide_data["title"]
            tf = slide.placeholders[1].text_frame
            for idx, bullet in enumerate(slide_data["bullets"]):
                if idx == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    
        # Bố cục 3: 2 Cột
        elif layout_id == 3:
            slide.shapes.title.text = slide_data["title"]
            
            # Cột trái (placeholder index 1)
            tf_left = slide.placeholders[1].text_frame
            for idx, bullet in enumerate(slide_data["column_1"]):
                if idx == 0:
                    tf_left.text = bullet
                else:
                    p = tf_left.add_paragraph()
                    p.text = bullet
                    
            # Cột phải (placeholder index 2)
            tf_right = slide.placeholders[2].text_frame
            for idx, bullet in enumerate(slide_data["column_2"]):
                if idx == 0:
                    tf_right.text = bullet
                else:
                    p = tf_right.add_paragraph()
                    p.text = bullet
                    
    output_path = "Demo_Dynamic.pptx"
    prs.save(output_path)
    print(f"✅ Thành công! Đã lưu file tại {output_path}")

if __name__ == "__main__":
    create_dynamic_presentation()
