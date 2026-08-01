import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

def extract_image_from_pdf(pdf_path, page_num, output_img):
    doc = fitz.open(pdf_path)
    page = doc[page_num]
    image_list = page.get_images(full=True)
    if image_list:
        largest_xref = None
        max_size = 0
        for img in image_list:
            xref = img[0]
            base_image = doc.extract_image(xref)
            size = len(base_image["image"])
            if size > max_size:
                max_size = size
                largest_xref = xref
                
        if largest_xref:
            base_image = doc.extract_image(largest_xref)
            image_bytes = base_image["image"]
            with open(output_img, "wb") as f:
                f.write(image_bytes)
            return True
    return False

def apply_color_block_theme(slide):
    # Nền xám nhạt (Clean aesthetic)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 249, 250)
    
    # Dải màu (Color Block) ở rìa trên
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0, 120, 215) # Corporate Blue
    shape1.line.fill.background()
    
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.15), Inches(3), Inches(0.1))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0, 180, 240) # Light Cyan
    shape2.line.fill.background()

def set_font(run, size_pt, bold=False, rgb=(50, 50, 50)):
    run.font.name = "Segoe UI"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*rgb)

def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    set_font(p.runs[0], 44, bold=True, rgb=(0, 60, 113)) # Dark Blue Title

def add_body(slide, text, top_inch=1.5):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top_inch), Inches(9), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    paragraphs = text.split('\n')
    for i, para_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para_text
        if p.runs:
            set_font(p.runs[0], 24, rgb=(60, 60, 60))

def main():
    pdf_path = "Chapter 11-Simple Linear Regression and Correlation.pdf"
    img_path = "temp_scatter.png"
    
    # Bước 1: Trích xuất biểu đồ ở trang 5 (index 4)
    has_image = extract_image_from_pdf(pdf_path, 4, img_path)
    
    # Dữ liệu nội dung thật tóm tắt từ PDF
    ai_json = [
        {
            "layout": "title_slide",
            "title": "Simple Linear Regression\nand Correlation",
            "subtitle": "Chapter 11 - Probability & Statistics"
        },
        {
            "layout": "one_column",
            "title": "Learning Objectives",
            "content": "1. Empirical Models\n2. Simple Linear Regression\n3. Properties of the Least Squares Estimators\n4. Hypothesis Tests in Simple Linear Regression\n5. Correlation"
        },
        {
            "layout": "one_column",
            "title": "Empirical Models",
            "content": "• Many problems in engineering and science involve exploring relationships between two or more variables.\n• Regression analysis is a statistical technique used to build models to predict outcomes.\n• Example: Predicting chemical product yield based on operating temperature."
        },
        {
            "layout": "image_slide",
            "title": "Scatter Diagram",
            "content": "Scatter Diagram of oxygen purity versus hydrocarbon level (Table 11-1).",
            "image_file": img_path if has_image else None
        },
        {
            "layout": "one_column",
            "title": "Simple Linear Regression Model",
            "content": "Based on the scatter diagram, we assume the mean of random variable Y is related to x by a straight-line relationship:\n\nE(Y|x) = β0 + β1x\n\nThe simple linear regression model with random error (ε) is given by:\n\nY = β0 + β1x + ε"
        }
    ]
    
    print("🚀 Bắt đầu tạo Slide bằng Native Code (Color Block)...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    for slide_data in ai_json:
        slide = prs.slides.add_slide(blank_layout)
        apply_color_block_theme(slide)
        
        if slide_data["layout"] == "title_slide":
            # Canh giữa trang
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data["title"]
            set_font(p.runs[0], 50, bold=True, rgb=(0, 60, 113))
            
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = slide_data["subtitle"]
            set_font(p2.runs[0], 28, rgb=(100, 100, 100))
            
        elif slide_data["layout"] == "one_column":
            add_title(slide, slide_data["title"])
            add_body(slide, slide_data["content"])
            
        elif slide_data["layout"] == "image_slide":
            add_title(slide, slide_data["title"])
            add_body(slide, slide_data["content"], top_inch=1.5)
            if slide_data.get("image_file"):
                # Chèn ảnh vào giữa màn hình
                slide.shapes.add_picture(slide_data["image_file"], Inches(1.5), Inches(2.2), width=Inches(7))
                
    output_path = "Slide_Hoan_Chinh_Color_Block.pptx"
    prs.save(output_path)
    
    if os.path.exists(img_path):
        os.remove(img_path)
        
    print(f"✅ Thành công! Đã lưu file tại {output_path}")

if __name__ == "__main__":
    main()
