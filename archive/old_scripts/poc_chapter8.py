import json
from pptx import Presentation

def create_dynamic_presentation():
    # Sử dụng Template từ Canva đã được chuẩn hóa (Chapter 8.pptx)
    template_path = "Chapter 8.pptx"
    prs = Presentation(template_path)
    
    # Xóa các slide mẫu cũ trong Chapter 8 (nếu có) để bắt đầu từ trang trắng
    # Lặp ngược để xóa an toàn
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    ai_director_plan = [
        {
            "layout_id": 0, # Title Slide
            "title": "Chương 8: Phân tích Dữ liệu Phức tạp",
            "subtitle": "Trình bày đẹp mắt với Canva & AI"
        },
        {
            "layout_id": 1, # Title and Content (1 Cột)
            "title": "Mục tiêu bài học",
            "bullets": [
                "Hiểu cách AI bóc tách cấu trúc Layout.",
                "Biết cách đẩy nội dung vào placeholder được định sẵn.",
                "Trải nghiệm sự khác biệt giữa text box tĩnh và layout động."
            ]
        },
        {
            "layout_id": 3, # Two Content (2 Cột)
            "title": "So sánh 2 giải pháp",
            "column_1": [
                "Cách cũ (Cứng nhắc)",
                "Nhồi nhét tất cả vào 1 slide.",
                "Chữ bị vỡ, tràn ra ngoài khung.",
                "Chỉ dùng được 1 mẫu duy nhất."
            ],
            "column_2": [
                "Cách mới (AI Director)",
                "Tự động chia tách nếu dài.",
                "Giữ nguyên 100% đồ họa Canva.",
                "Đổi Layout liên tục cho bớt nhàm chán."
            ]
        }
    ]
    
    print("🚀 Bắt đầu tạo Slide từ mẫu Chapter 8.pptx...")
    
    for slide_data in ai_director_plan:
        layout_id = slide_data["layout_id"]
        slide_layout = prs.slide_layouts[layout_id]
        slide = prs.slides.add_slide(slide_layout)
        
        # Bố cục 0: Title Slide
        if layout_id == 0:
            slide.shapes.title.text = slide_data["title"]
            slide.placeholders[1].text = slide_data["subtitle"]
            
        # Bố cục 1: 1 Cột
        elif layout_id == 1:
            slide.shapes.title.text = slide_data["title"]
            tf = slide.placeholders[1].text_frame
            for idx, bullet in enumerate(slide_data["bullets"]):
                if idx == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    
        # Bố cục 3: 2 Cột
        elif layout_id == 3:
            slide.shapes.title.text = slide_data["title"]
            
            # Cột trái (placeholder index 1)
            tf_left = slide.placeholders[1].text_frame
            for idx, bullet in enumerate(slide_data["column_1"]):
                if idx == 0:
                    tf_left.text = bullet
                else:
                    p = tf_left.add_paragraph()
                    p.text = bullet
                    
            # Cột phải (placeholder index 2)
            tf_right = slide.placeholders[2].text_frame
            for idx, bullet in enumerate(slide_data["column_2"]):
                if idx == 0:
                    tf_right.text = bullet
                else:
                    p = tf_right.add_paragraph()
                    p.text = bullet
                    
    output_path = "Slide_Hoan_Chinh_Chapter_8.pptx"
    prs.save(output_path)
    print(f"✅ Thành công! Đã lưu file tại {output_path}")

if __name__ == "__main__":
    create_dynamic_presentation()
