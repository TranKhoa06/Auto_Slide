from pptx import Presentation

def prepare_template(input_path, output_path):
    prs = Presentation(input_path)
    
    # Bỏ qua slide 0 (thường là tiêu đề) và slide 1 (thành viên nhóm)
    # Chúng ta sẽ chỉnh sửa từ Slide số 2 trở đi để làm khung {{TITLE}} và {{CONTENT}}
    for i in range(2, len(prs.slides)):
        slide = prs.slides[i]
        
        # Tìm các khung text
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text: # Nếu có chữ
                    text_shapes.append((shape, len(text), text))
        
        # Sắp xếp các khung chữ theo độ dài nội dung giảm dần
        text_shapes.sort(key=lambda x: x[1], reverse=True)
        
        if len(text_shapes) >= 2:
            # Khung chữ dài nhất -> {{CONTENT}}
            content_shape = text_shapes[0][0]
            # Khung chữ dài thứ 2 (thường là tiêu đề nhỏ hơn) -> {{TITLE}}
            title_shape = text_shapes[1][0]
            
            # Thay thế chữ nhưng GIỮ NGUYÊN ĐỊNH DẠNG (Bằng cách chỉ đổi Run đầu tiên và xóa các Run sau)
            def replace_text_preserve_format(shape, new_text):
                for paragraph in shape.text_frame.paragraphs:
                    if len(paragraph.runs) > 0:
                        paragraph.runs[0].text = new_text
                        # Xóa các run còn lại để không bị lặp chữ
                        for r in paragraph.runs[1:]:
                            r.text = ""
                            
            replace_text_preserve_format(content_shape, "{{CONTENT}}")
            replace_text_preserve_format(title_shape, "{{TITLE}}")
            
        elif len(text_shapes) == 1:
            replace_text_preserve_format(text_shapes[0][0], "{{CONTENT}}")

    prs.save(output_path)
    print(f"Đã tạo thành công template: {output_path}")

if __name__ == "__main__":
    prepare_template("BAN_MAU.pptx", "BAN_MAU_TEMPLATE.pptx")
